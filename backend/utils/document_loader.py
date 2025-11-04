import fitz  # PyMuPDF
import pptx
import openpyxl
from PIL import Image
import pytesseract
import io
import numpy as np
import cv2
from docx import Document


# Optional: specify Tesseract path (for Windows)
#pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


def extract_text_from_file(path: str, filename: str) -> str:
    ext = filename.split(".")[-1].lower()

    try:
        if ext == "pdf":
            return extract_from_pdf(path)
        elif ext == "docx":
            return extract_from_docx(path)
        elif ext == "txt":
            return open(path, "r", encoding="utf-8", errors="ignore").read()
        elif ext == "xlsx":
            return extract_from_xlsx(path)
        elif ext == "pptx":
            return extract_from_pptx(path)
        else:
            return f"[Unsupported file type: {ext}]"
    except Exception as e:
        return f"[Error reading {filename}: {str(e)}]"


# ---------------- PDF Extractor with OCR ---------------- #
def extract_from_pdf(path: str) -> str:
    """
    Extract text from PDFs, with OCR fallback for scanned/image-only pages.
    """
    text = ""
    with fitz.open(path) as doc:
        for page_index, page in enumerate(doc):
            page_text = page.get_text("text")
            if page_text.strip():
                # Normal text extraction
                text += page_text + "\n"
            else:
                # OCR fallback for image-based pages
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # upscale for better OCR accuracy
                img_data = pix.tobytes("png")
                image = Image.open(io.BytesIO(img_data))

                # Convert to grayscale and threshold to improve OCR
                gray = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2GRAY)
                _, thresh = cv2.threshold(gray, 150, 255, cv2.THRESH_BINARY)

                custom_config = r"--oem 3 --psm 6"
                ocr_text = pytesseract.image_to_string(thresh, config=custom_config)
                text += f"\n[Page {page_index + 1} OCR]\n" + ocr_text.strip() + "\n"

    return text


# ---------------- Word Extractor ---------------- #
def extract_from_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join([para.text for para in doc.paragraphs])


# ---------------- Excel Extractor ---------------- #
def extract_from_xlsx(path: str) -> str:
    wb = openpyxl.load_workbook(path, data_only=True)
    text = ""
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            text += " ".join([str(cell) for cell in row if cell]) + "\n"
    return text


# ---------------- PowerPoint Extractor ---------------- #
def extract_from_pptx(path: str) -> str:
    pres = pptx.Presentation(path)
    text = ""
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
